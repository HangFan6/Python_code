# -*- coding:utf-8 -*-
"""
作者：HET
日期：2023年04月14日
"""
import pptx

p=pptx.Presentation()  # 生成PPT对象
layout=p.slide_layouts[1] # 选择布局
# 0 title
# 1 title content
slide=p.slides.add_slide(layout)
p.save('test1.pptx')


