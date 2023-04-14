# -*- coding:utf-8 -*-
"""
作者：HET
日期：2023年04月14日
"""
import pptx

p=pptx.Presentation('test2.pptx')
# 获取每张幻灯片
for slide in p.slides:
    # 获取幻灯片中的形状
    for shape in slide.shapes:
        # 判断是否为文本类型
        if shape.has_text_frame:
            print(shape.text_frame.text)  # 获取文本内容
        if shape.has_table:
            for cell in shape.table.iter_cells():
                print(cell.text)  # 读取表格内容


