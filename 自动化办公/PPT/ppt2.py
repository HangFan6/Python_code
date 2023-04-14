# -*- coding:utf-8 -*-
"""
作者：HET
日期：2023年04月14日
"""
import pptx
from pptx.util import Pt,Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

p=pptx.Presentation()
layout=p.slide_layouts[1]  # title content
slide=p.slides.add_slide(layout)
# 添加标题
title=slide.placeholders[0]   # 0 title
title.text='题目'
# 添加段落
placeholder=slide.placeholders[1]   # 1 content
# placeholder.text='欢迎学习PPT制作\n欢迎学习python'  # 添加段落内容
# 添加段落1
paragraph_1=placeholder.text_frame.add_paragraph()
paragraph_1.text='欢迎学习PPT制作'
paragraph_1.bold=True  # 字体加粗
paragraph_1.font.italic=True  # 字体倾斜
paragraph_1.font.size=Pt(16)  # 设置字号
paragraph_1.font.underline=True  # 添加下划线
paragraph_1.alignment=PP_PARAGRAPH_ALIGNMENT.CENTER  # 段落居中
# 添加段落2
paragraph_2=placeholder.text_frame.add_paragraph()
paragraph_2.text='欢迎学习python'
paragraph_2.font.size=Pt(25)  # 设置字号
paragraph_2.alignment=PP_PARAGRAPH_ALIGNMENT.RIGHT  # 段落居中
# 自定义段落
layout=p.slide_layouts[6]
slide=p.slides.add_slide(layout)
left=top=Inches(1)  # 定义 左、上 边距
width=Inches(5)
height=Inches(1)
box=slide.shapes.add_textbox(left,top,width,height)  # 定义文本框
para=box.text_frame.add_paragraph()  # 添加段落
para.text='今天又学到了新知识'
para.alignment=PP_PARAGRAPH_ALIGNMENT.CENTER
para.font.size=Pt(25)
para.font.color.rgb=RGBColor(12,125,0)
para.font.name='微软雅黑'

# 插入表格
layout=p.slide_layouts[1]
slide=p.slides.add_slide(layout)
rows=10
cols=2
left=top=Inches(2)
width=Inches(6.0)
height=Inches(1.0)
table=slide.shapes.add_table(rows,cols,left, top, width, height).table
for index,_ in enumerate(range(rows)):
    for sub_index in range(cols):
        table.cell(index,sub_index).text='%s:%s'%(index,sub_index)

# 插入图片
layout = p.slide_layouts[1]
slide = p.slides.add_slide(layout)
left=top=Inches(1)
width=Inches(6)
image = slide.shapes.add_picture('../test3/images/imooc_1.jpg', left, top, width)  # height工具宽度自动匹配

p.save('test2.pptx')


